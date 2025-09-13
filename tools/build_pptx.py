#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
將 docs/slides_outline_zh.md + docs/slides_notes_zh.md 轉為 PowerPoint 檔 (.pptx)

需求:
  pip install python-pptx

輸出:
  docs/edge_voice_kws_methodology.pptx
"""

from pathlib import Path
import re
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
try:
    from PIL import Image  # 可選，用來從校徽取代表色
except Exception:
    Image = None

ROOT = Path(__file__).resolve().parents[1]
OUTLINE = ROOT / "docs" / "slides_outline_zh.md"
NOTES = ROOT / "docs" / "slides_notes_zh.md"
METHODOLOGY = ROOT / "docs" / "methodology_zh.md"
PINOUT = ROOT / "pinout.jpg"
OUTPUT = ROOT / "docs" / "edge_voice_kws_methodology.pptx"
LOGO_PRIMARY = ROOT / "tools" / "ncue logo.png"
LOGO_FALLBACK = ROOT / "tools" / "logo.png"
SELECTED_LOGO = None
DEFAULT_TITLE_RGB = RGBColor(0, 67, 137)  # 預設學院藍色系，若無法自動取色
TITLE_RGB = DEFAULT_TITLE_RGB

# 基本資訊（可自訂）
TITLE = "邊緣端語音指令系統"
SUBTITLE = "ESP32‑S3 + INMP441 + MQTT + 知識蒸餾"
AUTHOR = "國立彰化師範大學｜資訊工程學系 碩二｜M1354020 林昀佑"
ADVISORS = [
    "國立彰化師範大學 易昶霈 教授",
    "建國科大 沈慧宇 教授",
]
DATE_STR = datetime.now().strftime("%Y-%m-%d")
THEME_FONT = "Microsoft JhengHei"


def parse_outline(text: str):
    sections = []
    current = None
    for raw in text.splitlines():
        line = raw.strip()
        if not line:
            continue
        # 標題行: 以 數字. 開頭
        if re.match(r"^\d+\.", line):
            title = re.sub(r"^\d+\.\s*", "", line).strip()
            if current:
                sections.append(current)
            current = {"title": title, "bullets": []}
        elif line.startswith("- "):
            if current is None:
                current = {"title": "(未命名)", "bullets": []}
            current["bullets"].append(line[2:].strip())
    if current:
        sections.append(current)
    return sections


def parse_notes(text: str):
    # 以關鍵詞為索引
    mapping = {}
    key = None
    for raw in text.splitlines():
        line = raw.strip()
        if not line:
            continue
        if line.startswith("- "):
            body = line[2:].strip()
            # 嘗試抓出「關鍵詞：內容」
            if "：" in body:
                key, rest = body.split("：", 1)
                key = key.strip()
                mapping[key] = [rest.strip()] if rest.strip() else []
            else:
                if key:
                    mapping.setdefault(key, []).append(body)
        else:
            if key:
                mapping.setdefault(key, []).append(line)
    return mapping


def parse_methodology_sections(text: str):
    # 擷取以 "## <number>. <title>" 開頭的章節，收集其後的短行與條目
    sections = {}
    current = None
    for raw in text.splitlines():
        line = raw.strip()
        if not line:
            continue
        m = re.match(r"^##\s*\d+\.?\s*(.+)$", line)
        if m:
            current = m.group(1).strip()
            sections[current] = []
            continue
        if current:
            # 只收集簡短 bullet 或重要行
            if line.startswith("- "):
                sections[current].append(line[2:].strip())
            elif len(line) < 120 and not line.startswith("```"):
                # 短句也視為重點
                sections[current].append(line)
    return sections


def find_notes_for(title: str, notes_map: dict):
    # 以關鍵字模糊匹配
    keys = [
        ("封面", ["封面"]),
        ("動機", ["動機"]),
        ("架構", ["架構"]),
        ("邊緣", ["邊緣"]),
        ("MQTT", ["MQTT"]),
        ("蒸餾", ["蒸餾"]),
        ("部署", ["部署", "OTA"]),
        ("評估", ["評估"]),
        ("Demo", ["Demo"]),
        ("風險", ["風險"]),
        ("時程", ["時程"]),
    ]
    t = title.lower()
    for label, toks in keys:
        for tok in toks:
            if tok.lower() in t:
                if label in notes_map:
                    return notes_map[label]
    return []


def apply_text_style(text_frame, font_size=Pt(20), bold=False, color=RGBColor(20, 20, 20)):
    for p in text_frame.paragraphs:
        for r in p.runs:
            r.font.size = font_size
            r.font.bold = bold
            if THEME_FONT:
                r.font.name = THEME_FONT
            r.font.color.rgb = color


def set_theme_fonts(prs: Presentation):
    # 本函式目前不修改全域主題，改用逐段設定
    return prs


def add_footer(prs: Presentation, slide, page_num: int, total_pages: int):
    # 右下角頁碼與作者/日期
    left = prs.slide_width - Inches(3)
    top = prs.slide_height - Inches(0.6)
    width = Inches(2.8)
    height = Inches(0.4)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.clear()
    tf.text = f"{AUTHOR}  {DATE_STR}    {page_num}/{total_pages}"
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    apply_text_style(tf, font_size=Pt(10), color=RGBColor(90, 90, 90))


def add_logo(prs: Presentation, slide, is_title: bool = False):
    global SELECTED_LOGO
    logo_path = SELECTED_LOGO
    if logo_path is None:
        if LOGO_PRIMARY.exists():
            logo_path = LOGO_PRIMARY
        elif LOGO_FALLBACK.exists():
            logo_path = LOGO_FALLBACK
        SELECTED_LOGO = logo_path
    if not logo_path:
        return
    # 右上角放置校徽
    margin = Inches(0.35 if not is_title else 0.4)
    width = Inches(0.9 if not is_title else 1.2)
    left = prs.slide_width - width - margin
    top = margin
    try:
        slide.shapes.add_picture(str(logo_path), left, top, width=width)
    except Exception:
        pass


def extract_logo_color(path: Path):
    if Image is None or not path or not path.exists():
        return None
    try:
        img = Image.open(path).convert("RGB")
        img = img.resize((64, 64))
        pixels = list(img.getdata())
        # 過濾掉近白、近黑
        filtered = []
        for (r, g, b) in pixels:
            if r > 235 and g > 235 and b > 235:
                continue
            if r < 10 and g < 10 and b < 10:
                continue
            filtered.append((r, g, b))
        if not filtered:
            filtered = pixels
        # 取平均色
        n = len(filtered)
        r = sum(p[0] for p in filtered) // n
        g = sum(p[1] for p in filtered) // n
        b = sum(p[2] for p in filtered) // n
        return RGBColor(r, g, b)
    except Exception:
        return None


def add_advisors_slide(prs: Presentation):
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "指導教授"
    box = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), prs.slide_width - Inches(2.4), Inches(2.5))
    tf = box.text_frame
    tf.clear()
    if ADVISORS:
        tf.text = ADVISORS[0]
        apply_text_style(tf, font_size=Pt(24))
        for adv in ADVISORS[1:]:
            p = tf.add_paragraph()
            p.text = adv
            p.level = 0
    return slide


def add_system_architecture_slide(prs: Presentation):
    # 用方塊與箭頭構建簡單架構圖
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "系統架構（端‑雲協作）"

    def box(x, y, w, h, text):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))  # 1=Rectangle
        shape.text = text
        apply_text_style(shape.text_frame, font_size=Pt(16), bold=True)
        return shape

    # 佈局
    edge = box(0.6, 1.8, 3.4, 2.0, "ESP32‑S3 + INMP441\nVAD / 特徵 (log‑Mel/MFCC)\n本地學生模型")
    broker = box(4.2, 1.0, 2.6, 1.2, "MQTT Broker")
    server = box(7.2, 1.8, 3.2, 2.0, "伺服器端\n教師模型/蒸餾/訓練\n特徵接收 + 推論回覆")
    ota = box(7.2, 4.2, 3.2, 1.2, "模型倉庫/OTA 發佈")

    # 用簡單線段表示資料流
    def arrow(x1, y1, x2, y2):
        slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))

    # edge → broker → server（特徵上行）
    arrow(3.9, 2.8, 4.2, 2.1)
    arrow(6.8, 2.1, 7.2, 2.8)
    # server → broker → edge（推論/控制下發）
    arrow(7.2, 2.3, 6.8, 1.6)
    arrow(4.2, 1.6, 3.9, 2.3)
    # server → OTA（模型發佈）
    arrow(8.8, 3.8, 8.8, 4.2)

    return slide


def add_image_slide(prs: Presentation, title: str, image_path: Path):
    if not image_path.exists():
        return None
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    pic_left = Inches(1)
    pic_top = Inches(1.5)
    pic_width = prs.slide_width - Inches(2)
    slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=pic_width)
    return slide


def main():
    outline_text = OUTLINE.read_text(encoding="utf-8")
    notes_text = NOTES.read_text(encoding="utf-8")
    method_text = METHODOLOGY.read_text(encoding="utf-8")
    sections = parse_outline(outline_text)
    notes_map = parse_notes(notes_text)

    prs = Presentation()
    set_theme_fonts(prs)

    # 嘗試從校徽取主要顏色
    global TITLE_RGB
    logo_candidate = LOGO_PRIMARY if LOGO_PRIMARY.exists() else (LOGO_FALLBACK if LOGO_FALLBACK.exists() else None)
    col = extract_logo_color(logo_candidate) if logo_candidate else None
    if col is not None:
        TITLE_RGB = col

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = TITLE
    slide.placeholders[1].text = SUBTITLE
    if AUTHOR:
        # 加作者/日期於副標題下方
        tx = slide.shapes.add_textbox(Inches(0.9), Inches(4.5), prs.slide_width - Inches(1.8), Inches(0.4))
        tf = tx.text_frame
        tf.text = f"{AUTHOR} — {DATE_STR}"
        apply_text_style(tf, font_size=Pt(14), color=RGBColor(70, 70, 70))
    if ADVISORS:
        txa = slide.shapes.add_textbox(Inches(0.9), Inches(5.0), prs.slide_width - Inches(1.8), Inches(0.6))
        tfa = txa.text_frame
        tfa.text = "指導教授：" + "；".join(ADVISORS)
        apply_text_style(tfa, font_size=Pt(14), color=RGBColor(70, 70, 70))
    # 附上簡短 notes
    ns = find_notes_for("封面", notes_map)
    if ns:
        slide.notes_slide.notes_text_frame.text = "\n".join(ns)
    add_logo(prs, slide, is_title=True)

    # 增加系統架構與 Pinout 圖片
    add_advisors_slide(prs)
    add_system_architecture_slide(prs)
    add_image_slide(prs, "ESP32‑S3 引腳圖（Pinout）", PINOUT)

    # 從 methodology 擴充內容
    method_sections = parse_methodology_sections(method_text)

    # Section slides
    bullet_layout = prs.slide_layouts[1]
    for sec in sections:
        s = prs.slides.add_slide(bullet_layout)
        s.shapes.title.text = sec["title"]
        tf = s.shapes.placeholders[1].text_frame
        tf.clear()
        first = True
        for b in sec["bullets"]:
            if first:
                tf.text = b
                first = False
            else:
                tf.add_paragraph().text = b
        # 若 methodology 有對應章節，補一頁詳細
        for key in method_sections.keys():
            if key in sec["title"]:
                s2 = prs.slides.add_slide(bullet_layout)
                s2.shapes.title.text = f"詳細 — {sec['title']}"
                tf2 = s2.shapes.placeholders[1].text_frame
                tf2.clear()
                for i, b2 in enumerate(method_sections[key][:8]):
                    if i == 0:
                        tf2.text = b2
                    else:
                        tf2.add_paragraph().text = b2
                break
        ns = find_notes_for(sec["title"], notes_map)
        if ns:
            s.notes_slide.notes_text_frame.text = "\n".join(ns)

    # Thank you / Q&A
    end_slide = prs.slides.add_slide(title_slide_layout)
    end_slide.shapes.title.text = "Q&A"
    end_slide.placeholders[1].text = "謝謝指教"

    # 套用樣式（字體/大小）並加入頁腳
    total = len(prs.slides)
    for idx, sl in enumerate(prs.slides, start=1):
        # 處理標題與內容字型大小
        for shape in sl.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            tf = shape.text_frame
            # 標題較大，其餘內容較小
            if shape == sl.shapes.title:
                apply_text_style(tf, font_size=Pt(40), bold=True, color=TITLE_RGB)
            else:
                apply_text_style(tf, font_size=Pt(20))
        add_footer(prs, sl, idx, total)
        # 為非標題頁加校徽
        if idx != 1:
            add_logo(prs, sl, is_title=False)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"✅ 產生完成: {OUTPUT}")


if __name__ == "__main__":
    main()
